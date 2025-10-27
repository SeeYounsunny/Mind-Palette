#!/usr/bin/env python3
"""
Mind Palette - 컬러 일기 Professional Presentation Generator
Creates a professional PowerPoint presentation with modern design
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

def hex_to_rgb(hex_color):
    """Convert hex color to RGB tuple"""
    hex_color = hex_color.lstrip('#')
    return tuple(int(hex_color[i:i+2], 16) for i in (0, 2, 4))

def add_gradient_background(slide, color1, color2):
    """Add a subtle gradient background to the slide"""
    background = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        0, 0,
        Inches(10), Inches(5.625)  # 16:9 와이드 형식
    )
    background.line.fill.background()
    fill = background.fill
    fill.gradient()
    fill.gradient_angle = 45
    fill.gradient_stops[0].color.rgb = RGBColor(*hex_to_rgb(color1))
    fill.gradient_stops[1].color.rgb = RGBColor(*hex_to_rgb(color2))
    fill.gradient_stops[0].color.brightness = 0.95
    fill.gradient_stops[1].color.brightness = 0.90
    slide.shapes._spTree.remove(background._element)
    slide.shapes._spTree.insert(2, background._element)

def add_title_with_style(slide, title_text, top, left, width, height, font_size=44, bold=True):
    """Add a styled title text box"""
    title_box = slide.shapes.add_textbox(left, top, width, height)
    text_frame = title_box.text_frame
    text_frame.word_wrap = True
    p = text_frame.paragraphs[0]
    p.text = title_text
    p.alignment = PP_ALIGN.CENTER
    run = p.runs[0]
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = RGBColor(*hex_to_rgb('#667eea'))
    run.font.name = 'Arial'
    return title_box

def add_text_with_style(slide, text, top, left, width, height, font_size=18, color='#333333', alignment=PP_ALIGN.LEFT):
    """Add a styled text box"""
    text_box = slide.shapes.add_textbox(left, top, width, height)
    text_frame = text_box.text_frame
    text_frame.word_wrap = True
    text_frame.vertical_anchor = MSO_ANCHOR.TOP
    p = text_frame.paragraphs[0]
    p.text = text
    p.alignment = alignment
    run = p.runs[0]
    run.font.size = Pt(font_size)
    run.font.color.rgb = RGBColor(*hex_to_rgb(color))
    run.font.name = 'Arial'
    return text_box

def add_bullet_points(slide, points, top, left, width, font_size=20):
    """Add bullet points to slide"""
    text_box = slide.shapes.add_textbox(left, top, width, Inches(4))
    text_frame = text_box.text_frame
    text_frame.word_wrap = True
    for i, point in enumerate(points):
        if i == 0:
            p = text_frame.paragraphs[0]
        else:
            p = text_frame.add_paragraph()
        p.text = point
        p.level = 0
        p.space_before = Pt(12) if i > 0 else Pt(0)
        run = p.runs[0]
        run.font.size = Pt(font_size)
        run.font.color.rgb = RGBColor(51, 51, 51)
        run.font.name = 'Arial'

def add_decorative_shape(slide, top, left, width, height, color):
    """Add a decorative shape element"""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        left, top, width, height
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(*hex_to_rgb(color))
    shape.line.fill.background()
    return shape

# Create presentation
prs = Presentation()
prs.slide_width = Inches(10)      # 16:9 와이드 형식
prs.slide_height = Inches(5.625)  # 16:9 와이드 형식

# Define colors
PRIMARY_COLOR = '#667eea'
SECONDARY_COLOR = '#764ba2'
ACCENT_COLOR = '#f093fb'

# ====================================================================
# SLIDE 1 - TITLE
# ====================================================================
print("Creating Slide 1 - Title...")
slide1_layout = prs.slide_layouts[6]  # Blank layout
slide1 = prs.slides.add_slide(slide1_layout)

add_gradient_background(slide1, '#ffffff', '#f8f9ff')
add_decorative_shape(slide1, Inches(0.2), Inches(0.2), Inches(1.5), Inches(0.05), PRIMARY_COLOR)
add_decorative_shape(slide1, Inches(4.5), Inches(8.5), Inches(0.4), Inches(0.4), SECONDARY_COLOR)

# Main title
add_title_with_style(slide1, 'Mind Palette', Inches(1.5), Inches(0.8), Inches(7), Inches(0.8), font_size=56, bold=True)

# Subtitle
subtitle = slide1.shapes.add_textbox(Inches(1.5), Inches(2.1), Inches(7), Inches(0.6))
tf = subtitle.text_frame
p = tf.paragraphs[0]
p.text = '색채를 통해 나를 이해하는 감정 리플렉션 도구'
p.alignment = PP_ALIGN.CENTER
run = p.runs[0]
run.font.size = Pt(28)
run.font.color.rgb = RGBColor(*hex_to_rgb(SECONDARY_COLOR))
run.font.name = 'Arial'

# Description
desc = slide1.shapes.add_textbox(Inches(2), Inches(2.8), Inches(6), Inches(0.5))
tf = desc.text_frame
p = tf.paragraphs[0]
p.text = '33색 팔레트로 기록하는 나의 감정 여정'
p.alignment = PP_ALIGN.CENTER
run = p.runs[0]
run.font.size = Pt(18)
run.font.color.rgb = RGBColor(102, 102, 102)
run.font.name = 'Arial'

# Team info
team_info = slide1.shapes.add_textbox(Inches(2.5), Inches(3.6), Inches(5), Inches(1.5))
tf = team_info.text_frame
tf.word_wrap = True

info_lines = [
    '👥 팀원: 5명 (PM, 디자이너, 개발자, AI/데이터, 발표)',
    '⏱️ 개발 기간: 1주 MVP 집중 개발',
    '🎯 목표: 컬러 기반 감정 트래킹 앱 완성'
]

for i, line in enumerate(info_lines):
    if i == 0:
        p = tf.paragraphs[0]
    else:
        p = tf.add_paragraph()
    p.text = line
    p.alignment = PP_ALIGN.CENTER
    p.space_before = Pt(8) if i > 0 else Pt(0)
    run = p.runs[0]
    run.font.size = Pt(18)
    run.font.color.rgb = RGBColor(68, 68, 68)
    run.font.name = 'Arial'

# ====================================================================
# SLIDE 2 - PROBLEM STATEMENT
# ====================================================================
print("Creating Slide 2 - Problem Statement...")
slide2_layout = prs.slide_layouts[6]
slide2 = prs.slides.add_slide(slide2_layout)

add_gradient_background(slide2, '#ffffff', '#fff8fb')
add_decorative_shape(slide2, Inches(0.2), Inches(0.2), Inches(1.5), Inches(0.05), PRIMARY_COLOR)

# Title
add_title_with_style(slide2, '감정을 기록하고 표현하는 어려움', Inches(0.5), Inches(0.6), Inches(9), Inches(0.7), font_size=36)

# Subtitle
add_text_with_style(slide2, '현재의 문제점', Inches(1.5), Inches(1.3), Inches(7), Inches(0.4), font_size=22, color=SECONDARY_COLOR)

# Problem points
problems = [
    '💭 단순 텍스트로는 표현하기 어려운 복잡한 감정의 기록',
    '📊 감정의 변화 패턴을 체계적으로 파악하기 어려움',
    '💔 감정 기록의 지속성 부족으로 중단되는 경우가 많음',
    '🎨 자신의 감정 상태를 시각적으로 인식하고 공유할 수 있는 도구 부족'
]

add_bullet_points(slide2, problems, Inches(1.9), Inches(1.5), Inches(7), font_size=16)

# Target users box
target_box = slide2.shapes.add_textbox(Inches(1.5), Inches(4.1), Inches(7), Inches(1.0))
tf = target_box.text_frame
tf.word_wrap = True

p = tf.paragraphs[0]
p.text = '타겟 사용자'
p.alignment = PP_ALIGN.LEFT
run = p.runs[0]
run.font.size = Pt(22)
run.font.bold = True
run.font.color.rgb = RGBColor(*hex_to_rgb(PRIMARY_COLOR))

p2 = tf.add_paragraph()
p2.text = '감정을 기록하고 싶은 모든 사람, 자기 이해를 원하는 사람, 감정 패턴을 발견하고 싶은 사람'
p2.space_before = Pt(8)
run2 = p2.runs[0]
run2.font.size = Pt(16)
run2.font.color.rgb = RGBColor(68, 68, 68)

# ====================================================================
# SLIDE 3 - SOLUTION
# ====================================================================
print("Creating Slide 3 - Solution...")
slide3_layout = prs.slide_layouts[6]
slide3 = prs.slides.add_slide(slide3_layout)

add_gradient_background(slide3, '#ffffff', '#f8f9ff')
add_decorative_shape(slide3, Inches(0.2), Inches(0.2), Inches(1.5), Inches(0.05), PRIMARY_COLOR)

# Title
add_title_with_style(slide3, 'Mind Palette', Inches(0.5), Inches(0.6), Inches(9), Inches(0.6), font_size=40)

# Subtitle
add_text_with_style(slide3, '컬러 기반 감정 트래킹 솔루션', Inches(1.2), Inches(1.3), Inches(7), Inches(0.4), font_size=22, color=SECONDARY_COLOR, alignment=PP_ALIGN.CENTER)

# Core features
features = [
    '🎨 33색 팔레트로 직관적인 감정 표현',
    '📅 캘린더 뷰로 시간에 따른 감정 패턴 시각화',
    '📊 기간별 감정 분석 및 인사이트 제공',
    '💾 LocalStorage 기반 안전한 데이터 저장',
    '📸 공유용 감정 팔레트 이미지 자동 생성'
]

add_bullet_points(slide3, features, Inches(1.8), Inches(1.5), Inches(7), font_size=15)

# Key values section
value_box = slide3.shapes.add_textbox(Inches(1.5), Inches(3.8), Inches(7), Inches(1.5))
tf = value_box.text_frame
tf.word_wrap = True

p = tf.paragraphs[0]
p.text = '핵심 가치'
p.alignment = PP_ALIGN.LEFT
run = p.runs[0]
run.font.size = Pt(22)
run.font.bold = True
run.font.color.rgb = RGBColor(*hex_to_rgb(PRIMARY_COLOR))

value_points = [
    '✨ 색상 선택만으로 쉽고 빠른 감정 기록',
    '🎯 개인 맞춤형 감정 팔레트 생성',
    '🔒 로컬 저장으로 개인정보 보호 완벽'
]

for point in value_points:
    p = tf.add_paragraph()
    p.text = point
    p.space_before = Pt(6)
    run = p.runs[0]
    run.font.size = Pt(15)
    run.font.color.rgb = RGBColor(68, 68, 68)

# ====================================================================
# SLIDE 4 - DEMO
# ====================================================================
print("Creating Slide 4 - Demo...")
slide4_layout = prs.slide_layouts[6]
slide4 = prs.slides.add_slide(slide4_layout)

add_gradient_background(slide4, '#ffffff', '#fff8fb')
add_decorative_shape(slide4, Inches(0.2), Inches(0.2), Inches(1.5), Inches(0.05), PRIMARY_COLOR)

# Title
add_title_with_style(slide4, 'Mind Palette 사용 방법', Inches(0.5), Inches(0.6), Inches(9), Inches(0.6), font_size=36)

# Demo scenario
add_text_with_style(slide4, '일기 작성 프로세스', Inches(1.5), Inches(1.3), Inches(7), Inches(0.4), font_size=22, color=SECONDARY_COLOR)

demo_steps = [
    '1️⃣ 33색 팔레트에서 지금 가장 끌리는 컬러 선택',
    '2️⃣ 감정 선택 (21가지) + 강도 조절 (1-5)',
    '3️⃣ 오늘의 에피소드를 자유롭게 기록',
    '4️⃣ 날씨와 날씨에 대한 느낌 선택',
    '5️⃣ 저장 후 캘린더에서 색상으로 확인'
]

add_bullet_points(slide4, demo_steps, Inches(1.8), Inches(1.5), Inches(7), font_size=15)

# Additional features
features_box = slide4.shapes.add_textbox(Inches(1.5), Inches(3.8), Inches(7), Inches(1.5))
tf = features_box.text_frame
tf.word_wrap = True

p = tf.paragraphs[0]
p.text = '추가 기능'
p.alignment = PP_ALIGN.LEFT
run = p.runs[0]
run.font.size = Pt(22)
run.font.bold = True
run.font.color.rgb = RGBColor(*hex_to_rgb(PRIMARY_COLOR))

feature_points = [
    '📊 감정 분석: 1주/1개월/3개월 단위 감정 패턴 분석',
    '📸 공유 이미지: 월간 감정 팔레트 이미지 자동 생성',
    '💾 데이터 내보내기: CSV 파일로 전체 데이터 백업'
]

for point in feature_points:
    p = tf.add_paragraph()
    p.text = point
    p.space_before = Pt(6)
    run = p.runs[0]
    run.font.size = Pt(15)
    run.font.color.rgb = RGBColor(68, 68, 68)

# ====================================================================
# SLIDE 5 - VALUE PROPOSITION
# ====================================================================
print("Creating Slide 5 - Value Proposition...")
slide5_layout = prs.slide_layouts[6]
slide5 = prs.slides.add_slide(slide5_layout)

add_gradient_background(slide5, '#ffffff', '#f8f9ff')
add_decorative_shape(slide5, Inches(0.2), Inches(0.2), Inches(1.5), Inches(0.05), PRIMARY_COLOR)

# Title
add_title_with_style(slide5, 'Mind Palette가 만들어내는 가치', Inches(0.5), Inches(0.5), Inches(9), Inches(0.6), font_size=34)

# User value
user_value_box = slide5.shapes.add_textbox(Inches(0.8), Inches(1.3), Inches(4), Inches(1.8))
tf = user_value_box.text_frame
tf.word_wrap = True

p = tf.paragraphs[0]
p.text = '사용자 가치'
run = p.runs[0]
run.font.size = Pt(20)
run.font.bold = True
run.font.color.rgb = RGBColor(*hex_to_rgb(PRIMARY_COLOR))

user_values = [
    '🎨 직관적인 색상 선택으로 쉬운 감정 기록',
    '📊 시간에 따른 감정 패턴의 시각적 발견',
    '💭 자신의 감정에 대한 깊은 이해',
    '📱 언제 어디서나 간편한 사용'
]

for value in user_values:
    p = tf.add_paragraph()
    p.text = value
    p.space_before = Pt(8)
    run = p.runs[0]
    run.font.size = Pt(14)
    run.font.color.rgb = RGBColor(68, 68, 68)

# Technical value
tech_box = slide5.shapes.add_textbox(Inches(5.2), Inches(1.3), Inches(4), Inches(1.8))
tf = tech_box.text_frame
tf.word_wrap = True

p = tf.paragraphs[0]
p.text = '기술적 특징'
run = p.runs[0]
run.font.size = Pt(20)
run.font.bold = True
run.font.color.rgb = RGBColor(*hex_to_rgb(SECONDARY_COLOR))

tech_values = [
    '⚡ React + Vite 기반 빠른 로딩',
    '💾 LocalStorage 완전 오프라인 지원',
    '🎨 반응형 디자인 (모바일 최적화)',
    '🔒 개인정보 보호 (서버 전송 없음)'
]

for value in tech_values:
    p = tf.add_paragraph()
    p.text = value
    p.space_before = Pt(8)
    run = p.runs[0]
    run.font.size = Pt(14)
    run.font.color.rgb = RGBColor(68, 68, 68)

# Technical achievements
achieve_box = slide5.shapes.add_textbox(Inches(0.8), Inches(3.3), Inches(4), Inches(1.2))
tf = achieve_box.text_frame
tf.word_wrap = True

p = tf.paragraphs[0]
p.text = '개발 성과'
run = p.runs[0]
run.font.size = Pt(16)
run.font.bold = True
run.font.color.rgb = RGBColor(*hex_to_rgb(PRIMARY_COLOR))

p = tf.add_paragraph()
p.text = '1주일 MVP 완성, React 기반 모던 UI, 감정 분석 기능 구현'
p.space_before = Pt(4)
run = p.runs[0]
run.font.size = Pt(12)
run.font.color.rgb = RGBColor(68, 68, 68)

# Future roadmap
roadmap_box = slide5.shapes.add_textbox(Inches(5.2), Inches(3.3), Inches(4), Inches(1.2))
tf = roadmap_box.text_frame
tf.word_wrap = True

p = tf.paragraphs[0]
p.text = '향후 계획'
run = p.runs[0]
run.font.size = Pt(18)
run.font.bold = True
run.font.color.rgb = RGBColor(*hex_to_rgb(SECONDARY_COLOR))

p = tf.add_paragraph()
p.text = 'PWA 지원, AI 분석 고도화, 모바일 앱, 클라우드 백업'
p.space_before = Pt(6)
run = p.runs[0]
run.font.size = Pt(13)
run.font.color.rgb = RGBColor(68, 68, 68)

# Closing statement
closing_box = slide5.shapes.add_textbox(Inches(1.5), Inches(4.8), Inches(7), Inches(0.6))
tf = closing_box.text_frame

p = tf.paragraphs[0]
p.text = '색채로 기록하는 나의 감정 여정'
p.alignment = PP_ALIGN.CENTER
run = p.runs[0]
run.font.size = Pt(20)
run.font.bold = True
run.font.color.rgb = RGBColor(*hex_to_rgb(PRIMARY_COLOR))

p2 = tf.add_paragraph()
p2.text = 'Mind Palette'
p2.alignment = PP_ALIGN.CENTER
run2 = p2.runs[0]
run2.font.size = Pt(24)
run2.font.bold = True
run2.font.color.rgb = RGBColor(*hex_to_rgb(SECONDARY_COLOR))

# Save presentation
output_file = 'presentation/Mind_Palette_Presentation.pptx'
prs.save(output_file)
print(f"\n✅ Presentation saved successfully: {output_file}")
print(f"📊 Total slides: {len(prs.slides)}")
print(f"📐 Slide dimensions: {prs.slide_width.inches}\" x {prs.slide_height.inches}\"")
